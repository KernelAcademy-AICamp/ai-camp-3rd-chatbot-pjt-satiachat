from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Teal & Coral 컬러 팔레트 (프로젝트 테마 컬러)
PRIMARY = RGBColor(0x5E, 0xA8, 0xA7)  # Teal
PRIMARY_DARK = RGBColor(0x27, 0x78, 0x84)
ACCENT = RGBColor(0xFE, 0x44, 0x47)  # Coral
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1C, 0x28, 0x33)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)

# Create presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """타이틀 슬라이드 추가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_DARK
    bg.line.fill.background()

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xA0, 0xD0, 0xCF)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title):
    """섹션 구분 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 좌측 컬러 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK

    return slide

def add_content_slide(prs, title, contents, two_column=False):
    """콘텐츠 슬라이드 (1컬럼 또는 2컬럼)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK

    if two_column:
        # 2컬럼 레이아웃
        left_contents, right_contents = contents

        # 왼쪽 컬럼
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_contents):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = DARK
            p.space_after = Pt(12)

        # 오른쪽 컬럼
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_contents):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = DARK
            p.space_after = Pt(12)
    else:
        # 1컬럼 레이아웃
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(contents):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = DARK
            p.space_after = Pt(14)

    return slide

def add_architecture_slide(prs):
    """아키텍처 다이어그램 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK

    # React App 박스
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(2.8), Inches(0.8))
    box1.fill.solid()
    box1.fill.fore_color.rgb = PRIMARY
    box1.line.color.rgb = PRIMARY_DARK
    tf = box1.text_frame
    tf.paragraphs[0].text = "React Frontend"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 화살표 (텍스트로 표현)
    arrow1 = slide.shapes.add_textbox(Inches(4.5), Inches(2.5), Inches(1), Inches(0.5))
    arrow1.text_frame.paragraphs[0].text = "↓"
    arrow1.text_frame.paragraphs[0].font.size = Pt(24)
    arrow1.text_frame.paragraphs[0].font.color.rgb = GRAY

    arrow2 = slide.shapes.add_textbox(Inches(6.5), Inches(2.5), Inches(1), Inches(0.5))
    arrow2.text_frame.paragraphs[0].text = "↓"
    arrow2.text_frame.paragraphs[0].font.size = Pt(24)
    arrow2.text_frame.paragraphs[0].font.color.rgb = GRAY

    arrow3 = slide.shapes.add_textbox(Inches(8.2), Inches(2.5), Inches(1), Inches(0.5))
    arrow3.text_frame.paragraphs[0].text = "↓"
    arrow3.text_frame.paragraphs[0].font.size = Pt(24)
    arrow3.text_frame.paragraphs[0].font.color.rgb = GRAY

    # Supabase Auth
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.3), Inches(2.5), Inches(1.2))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(0x3E, 0xCF, 0x8E)  # Supabase green
    box2.line.color.rgb = RGBColor(0x24, 0xB4, 0x7E)
    tf = box2.text_frame
    tf.paragraphs[0].text = "Supabase Auth"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Login / Signup"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Supabase DB
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(3.3), Inches(2.5), Inches(1.2))
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(0x3E, 0xCF, 0x8E)
    box3.line.color.rgb = RGBColor(0x24, 0xB4, 0x7E)
    tf = box3.text_frame
    tf.paragraphs[0].text = "Supabase DB"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "PostgreSQL + RLS"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # FastAPI
    box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(3.3), Inches(2.5), Inches(1.2))
    box4.fill.solid()
    box4.fill.fore_color.rgb = RGBColor(0x00, 0x96, 0x88)  # FastAPI teal
    box4.line.color.rgb = RGBColor(0x00, 0x7A, 0x6E)
    tf = box4.text_frame
    tf.paragraphs[0].text = "FastAPI Server"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "AI Chat + RAG"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 화살표
    arrow4 = slide.shapes.add_textbox(Inches(8.5), Inches(4.7), Inches(1), Inches(0.5))
    arrow4.text_frame.paragraphs[0].text = "↓"
    arrow4.text_frame.paragraphs[0].font.size = Pt(24)
    arrow4.text_frame.paragraphs[0].font.color.rgb = GRAY

    # OpenAI
    box5 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(5.3), Inches(2.5), Inches(1))
    box5.fill.solid()
    box5.fill.fore_color.rgb = RGBColor(0x10, 0xA3, 0x7F)  # OpenAI
    box5.line.color.rgb = RGBColor(0x0D, 0x8A, 0x6B)
    tf = box5.text_frame
    tf.paragraphs[0].text = "OpenAI API"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 설명 박스
    desc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.3), Inches(5), Inches(1.5))
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = LIGHT_BG
    desc_box.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    tf = desc_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Data Flow:"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK
    p = tf.add_paragraph()
    p.text = "• CRUD Operations → Supabase Direct"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p = tf.add_paragraph()
    p.text = "• AI Chatbot → FastAPI → OpenAI"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p = tf.add_paragraph()
    p.text = "• Auth → Supabase Auth (JWT)"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY

    return slide

def add_table_slide(prs, title, headers, rows):
    """테이블 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK

    # 테이블
    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5 * table_rows)).table

    # 헤더
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # 데이터
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK

    return slide

# ========== 슬라이드 생성 ==========

# 1. 타이틀 슬라이드
add_title_slide(prs, "DietRx Coach", "AI 기반 다이어트 코칭 플랫폼\nProject Technical Overview")

# 2. 프로젝트 개요
add_section_slide(prs, "01  Project Overview")

add_content_slide(prs, "프로젝트 소개", [
    "■ DietRx Coach",
    "   AI 페르소나 기반 다이어트 코칭 플랫폼",
    "",
    "■ 핵심 기능",
    "   • 식사 기록 및 칼로리 자동 계산 (한국 음식 11,000+ 데이터베이스)",
    "   • 체중/체성분 추적 및 시각화",
    "   • 주간 약물 복용 관리 (GLP-1 약물 특화: 위고비, 마운자로)",
    "   • AI 코칭 챗봇 (3가지 페르소나 선택 가능)",
    "   • 커뮤니티 게시판",
    "",
    "■ 타겟 사용자",
    "   • 체중 감량을 원하는 사용자",
    "   • GLP-1 계열 약물 복용자",
])

# 3. 기술 스택
add_section_slide(prs, "02  Technology Stack")

add_content_slide(prs, "Frontend Stack", ([
    "■ Core Framework",
    "   • React 18 + TypeScript",
    "   • Vite + SWC (빌드 도구)",
    "",
    "■ UI Components",
    "   • shadcn/ui (Radix Primitives)",
    "   • Tailwind CSS",
    "   • Lucide Icons",
    "",
    "■ State Management",
    "   • TanStack Query v5",
    "   • React Hook Form + Zod",
], [
    "■ Routing",
    "   • React Router DOM v6",
    "",
    "■ Charts & Visualization",
    "   • Recharts",
    "",
    "■ Date Handling",
    "   • date-fns",
    "",
    "■ Key Dependencies",
    "   • @supabase/supabase-js",
    "   • sonner (토스트)",
    "   • cmdk (커맨드 팔레트)",
]), two_column=True)

add_content_slide(prs, "Backend Stack", ([
    "■ API Server",
    "   • FastAPI (Python)",
    "   • Uvicorn (ASGI Server)",
    "   • Pydantic v2 (Validation)",
    "",
    "■ Database",
    "   • Supabase (PostgreSQL)",
    "   • Row Level Security (RLS)",
    "",
    "■ Authentication",
    "   • Supabase Auth",
    "   • JWT Token",
], [
    "■ AI / ML",
    "   • OpenAI GPT-4",
    "   • LlamaIndex (RAG)",
    "   • OpenAI Embeddings",
    "",
    "■ Infrastructure",
    "   • Supabase Cloud",
    "   • Edge Functions (Optional)",
    "",
    "■ Development",
    "   • Python 3.11+",
    "   • Node.js 18+",
]), two_column=True)

# 4. 아키텍처
add_section_slide(prs, "03  System Architecture")

add_architecture_slide(prs)

add_content_slide(prs, "Architecture Patterns", ([
    "■ Frontend Architecture",
    "   • Feature-based folder structure",
    "   • Custom hooks for data fetching",
    "   • Compound component patterns",
    "",
    "■ Data Access Pattern",
    "   • BaaS (Backend as a Service)",
    "   • Direct Supabase queries for CRUD",
    "   • FastAPI for AI/complex logic",
], [
    "■ Security",
    "   • Supabase RLS (Row Level Security)",
    "   • JWT-based authentication",
    "   • API key protection via server",
    "",
    "■ State Management",
    "   • Server state: TanStack Query",
    "   • Form state: React Hook Form",
    "   • Auth state: React Context",
]), two_column=True)

# 5. 주요 기능
add_section_slide(prs, "04  Key Features")

add_table_slide(prs, "Feature List",
    ["Feature", "Description", "Status"],
    [
        ["Dashboard", "오늘의 칼로리, 체중, 약물 요약", "✅ Complete"],
        ["Meals", "식사 기록, 음식 검색, 영양소 분석", "✅ Complete"],
        ["Medications", "주간 약물 관리, 복용 캘린더", "✅ Complete"],
        ["AI Chat", "페르소나 기반 AI 코칭", "✅ Complete"],
        ["RAG Chat", "약물 정보 RAG 기반 응답", "✅ Complete"],
        ["Board", "커뮤니티 게시판", "✅ Complete"],
        ["Settings", "프로필 설정, 테마", "✅ Complete"],
    ]
)

# 6. 데이터베이스 스키마
add_section_slide(prs, "05  Database Schema")

add_table_slide(prs, "Main Tables",
    ["Table", "Description", "Key Fields"],
    [
        ["user_profiles", "사용자 프로필", "nickname, weight, goal, persona"],
        ["meals", "식사 기록", "date, meal_type, total_calories"],
        ["meal_items", "식사 항목", "meal_id, food_name, calories, nutrients"],
        ["medications", "약물 정보", "name, dosage, dose_day, frequency"],
        ["medication_logs", "복용 기록", "medication_id, taken_at, status"],
        ["progress_logs", "체중 기록", "date, weight_kg, body_fat"],
        ["posts", "게시글", "tab, title, content, likes"],
        ["comments", "댓글", "post_id, content, author_id"],
    ]
)

# 7. 프로젝트 구조
add_section_slide(prs, "06  Project Structure")

add_content_slide(prs, "Directory Structure", ([
    "■ Frontend (src/)",
    "   ├── components/",
    "   │   ├── ui/          # shadcn components",
    "   │   ├── dashboard/   # Dashboard widgets",
    "   │   ├── meals/       # Meal components",
    "   │   ├── medications/ # Medication components",
    "   │   └── layout/      # App shell",
    "   ├── pages/           # Route pages",
    "   ├── hooks/           # Custom hooks",
    "   ├── lib/             # Utilities",
    "   ├── types/           # TypeScript types",
    "   └── contexts/        # React contexts",
], [
    "■ Backend (server/)",
    "   ├── api/",
    "   │   └── v1/          # API routes",
    "   ├── services/        # Business logic",
    "   ├── ai/              # AI integration",
    "   ├── rag/             # RAG components",
    "   ├── schemas/         # Pydantic models",
    "   └── core/            # Config, deps",
    "",
    "■ Database (supabase/)",
    "   ├── schema.sql       # DB schema",
    "   ├── migrations/      # Migrations",
    "   └── functions/       # Edge Functions",
]), two_column=True)

# 8. API 엔드포인트
add_section_slide(prs, "07  API Endpoints")

add_table_slide(prs, "FastAPI Endpoints",
    ["Method", "Endpoint", "Description"],
    [
        ["POST", "/api/v1/chat/message", "AI 챗봇 메시지 전송"],
        ["POST", "/api/v1/medication/ask", "약물 RAG 질문"],
        ["GET", "/api/v1/health/summary", "건강 데이터 요약"],
        ["GET", "/api/v1/meals/today", "오늘의 식사"],
        ["GET", "/api/v1/progress/weight", "체중 기록"],
    ]
)

# 9. 마무리
add_section_slide(prs, "08  Summary")

add_content_slide(prs, "Project Summary", [
    "■ 하이브리드 아키텍처",
    "   • Supabase BaaS + FastAPI AI Server",
    "   • 빠른 개발 속도와 AI 기능 확장성 확보",
    "",
    "■ 기술적 특징",
    "   • TypeScript 기반 타입 안전성",
    "   • TanStack Query로 서버 상태 최적화",
    "   • RAG 기반 AI 응답 품질 향상",
    "",
    "■ 향후 계획",
    "   • 알림 기능 (Push Notification)",
    "   • 리포트 자동 생성",
    "   • 모바일 앱 (React Native)",
])

# 마지막 슬라이드
add_title_slide(prs, "Thank You", "DietRx Coach - AI Diet Coaching Platform")

# 저장
prs.save('C:/workspace/mini_project/docs/DietRx_Coach_Overview.pptx')
print("Presentation saved: C:/workspace/mini_project/docs/DietRx_Coach_Overview.pptx")
